"""Text extraction — extract plain text from various document formats.

Supports optional OCR (Tesseract) and Vision (Claude) enrichment for
image-heavy pages during indexing. These are passed as service references
and are never called at query time.
"""

from __future__ import annotations

import io
import json
import logging
from typing import TYPE_CHECKING, Any

from gilbert.interfaces.knowledge import DocumentContent, DocumentType, ExtractionStats

if TYPE_CHECKING:
    from gilbert.core.services.ocr import OCRService
    from gilbert.core.services.vision import VisionService

logger = logging.getLogger(__name__)

# Sparse text threshold — pages with fewer chars than this are considered
# image-heavy and eligible for OCR/Vision enrichment.
_SPARSE_TEXT_THRESHOLD = 50

# Domain-specific prompt for PDF knowledge indexing. The vision
# backend's default prompt is general-purpose ("describe what's in
# this image") which is wrong for technical-document pages — we want
# strict information extraction from pinout tables, wiring diagrams,
# specs, etc., not natural-language scene narration. We pin our own
# prompt here so this caller's behavior is independent of how the
# operator tunes Settings → Vision → Prompt for other consumers
# (smart-glasses scene description, surveillance cameras, etc.).
_PDF_PAGE_VISION_PROMPT = (
    "Extract ALL technical content from this page image as plain "
    "structured text. Include: pinout tables, wiring diagrams, "
    "connector assignments, component specifications, part numbers, "
    "voltage/current ratings, communication protocols, dimensions, "
    "torque specs, and any other technical data. Reproduce tables as "
    "aligned text columns. Label diagram elements clearly (e.g., "
    "'Pin 1: CAN_H, Pin 2: CAN_L'). Do NOT describe the visual "
    "layout — extract the information content only. If the page "
    "contains no technical content, respond with an empty string."
)


def extract_text(
    content: DocumentContent,
    *,
    vision: VisionService | None = None,
    ocr: OCRService | None = None,
    max_chars: int = 200000,
) -> tuple[str, ExtractionStats]:
    """Extract plain text from a document.

    Args:
        content: Document content with metadata.
        vision: Optional VisionService for image-heavy page understanding.
        ocr: Optional OCRService for text extraction from images.
        max_chars: Maximum characters to extract (truncates beyond this).

    Returns:
        (extracted_text, stats) tuple.
    """
    stats = ExtractionStats()
    match content.meta.document_type:
        case DocumentType.TEXT | DocumentType.MARKDOWN | DocumentType.CSV:
            text = content.data.decode(content.encoding, errors="replace")
        case DocumentType.JSON:
            text = _extract_json(content.data, content.encoding)
        case DocumentType.YAML:
            text = _extract_yaml(content.data, content.encoding)
        case DocumentType.PDF:
            text = _extract_pdf(content.data, stats, vision=vision, ocr=ocr, max_chars=max_chars)
        case DocumentType.WORD:
            text = _extract_word(content.data, stats, ocr=ocr, vision=vision, max_chars=max_chars)
        case DocumentType.EXCEL:
            text = _extract_excel(content.data)
        case DocumentType.POWERPOINT:
            text = _extract_powerpoint(content.data)
        case _:
            text = content.data.decode(content.encoding, errors="replace")

    if len(text) > max_chars:
        text = text[:max_chars] + "\n... (truncated)"

    stats.total_chars = len(text)
    return text, stats


def _extract_json(data: bytes, encoding: str = "utf-8") -> str:
    """Pretty-print JSON for searchability."""
    try:
        parsed = json.loads(data.decode(encoding, errors="replace"))
        return json.dumps(parsed, indent=2, ensure_ascii=False)
    except (json.JSONDecodeError, ValueError):
        return data.decode(encoding, errors="replace")


def _extract_yaml(data: bytes, encoding: str = "utf-8") -> str:
    """Load and dump YAML as formatted text."""
    try:
        import yaml

        parsed = yaml.safe_load(data.decode(encoding, errors="replace"))
        return yaml.dump(parsed, default_flow_style=False, allow_unicode=True)
    except Exception:
        return data.decode(encoding, errors="replace")


def _extract_pdf(
    data: bytes,
    stats: ExtractionStats,
    *,
    vision: VisionService | None = None,
    ocr: OCRService | None = None,
    max_chars: int = 200000,
) -> str:
    """Extract text from PDF using PyMuPDF, with OCR and Vision enrichment.

    For each page:
    - Extract embedded text
    - If sparse text + images: run OCR (tesseract) at 300dpi
    - If sparse text + images + vision available: run Claude Vision at 200dpi
    """
    try:
        import pymupdf

        doc: Any = pymupdf.open(stream=data, filetype="pdf")  # type: ignore[no-untyped-call]
        text_parts: list[str] = []
        total_chars = 0
        stats.pages = len(doc)

        for page_num in range(len(doc)):
            page = doc[page_num]
            page_text: str = page.get_text()

            # Check for images and sparse text
            page_images = page.get_images()
            has_images = bool(page_images)
            sparse_text = len(page_text.strip()) < _SPARSE_TEXT_THRESHOLD

            if has_images:
                stats.images_found += len(page_images)

            # OCR: extract text from image-heavy or sparse pages
            if ocr and ocr.available and (sparse_text or has_images):
                try:
                    pix = page.get_pixmap(dpi=300)
                    ocr_text = _run_ocr_sync(ocr, pix)
                    if ocr_text:
                        if sparse_text:
                            page_text = f"[OCR] {ocr_text}"
                        else:
                            page_text = f"{page_text}\n[OCR] {ocr_text}"
                        stats.ocr_pages += 1
                        stats.ocr_chars += len(ocr_text)
                        logger.debug("OCR page %d: %d chars", page_num + 1, len(ocr_text))
                except Exception:
                    logger.warning("OCR failed on page %d", page_num + 1, exc_info=True)
                    stats.warnings.append(f"OCR failed on page {page_num + 1}")

            # Vision: semantic understanding of image-heavy pages
            if vision and vision.available and has_images and sparse_text:
                try:
                    pix_v = page.get_pixmap(dpi=200)
                    png_bytes = pix_v.tobytes("png")
                    vision_text = _run_vision_sync(vision, png_bytes)
                    if vision_text:
                        page_text = f"{page_text}\n[Vision] {vision_text}"
                        stats.vision_pages += 1
                        stats.vision_chars += len(vision_text)
                        logger.info(
                            "Vision page %d of %s: %d chars",
                            page_num + 1,
                            doc.name or "PDF",
                            len(vision_text),
                        )
                except Exception:
                    logger.warning("Vision failed on page %d", page_num + 1, exc_info=True)
                    stats.warnings.append(f"Vision failed on page {page_num + 1}")

            text_parts.append(f"[Page {page_num + 1}]\n{page_text}")
            total_chars += len(page_text)
            if total_chars > max_chars:
                text_parts.append("\n... (truncated — document continues)")
                break

        doc.close()
        return "\n".join(text_parts)
    except Exception:
        logger.warning("Failed to extract text from PDF", exc_info=True)
        return ""


def _run_ocr_sync(ocr: OCRService, pixmap: Any) -> str:
    """Run OCR on a PyMuPDF pixmap (synchronous — called from thread)."""
    import pytesseract  # type: ignore[import-untyped]
    from PIL import Image

    img = Image.frombytes("RGB", (pixmap.width, pixmap.height), pixmap.samples)
    return str(pytesseract.image_to_string(img)).strip()


def _run_vision_sync(vision: VisionService, png_bytes: bytes) -> str:
    """Run Vision on image bytes (synchronous — blocks, used during indexing only)."""
    import asyncio

    # Vision.describe_image is async, but during indexing we're already
    # in a thread via asyncio.to_thread, so we run a new event loop.
    try:
        loop = asyncio.get_event_loop()
        if loop.is_running():
            # We're in a thread called from the async event loop —
            # can't use loop.run_until_complete, use asyncio.run in a new loop
            import concurrent.futures

            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                future = pool.submit(
                    asyncio.run,
                    vision.describe_image(
                        png_bytes, "image/png", prompt=_PDF_PAGE_VISION_PROMPT
                    ),
                )
                return future.result(timeout=60)
        else:
            return loop.run_until_complete(
                vision.describe_image(
                    png_bytes, "image/png", prompt=_PDF_PAGE_VISION_PROMPT
                )
            )
    except Exception:
        # Simplest fallback: just run it directly
        return asyncio.run(
            vision.describe_image(
                png_bytes, "image/png", prompt=_PDF_PAGE_VISION_PROMPT
            )
        )


def _extract_word(
    data: bytes,
    stats: ExtractionStats,
    *,
    ocr: OCRService | None = None,
    vision: VisionService | None = None,
    max_chars: int = 200000,
) -> str:
    """Extract text from Word documents using python-docx."""
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        parts: list[str] = []
        total_chars = 0

        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
                total_chars += len(para.text)

        # Extract table content
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                if cells:
                    line = " | ".join(cells)
                    parts.append(line)
                    total_chars += len(line)

        stats.pages = 1

        # OCR embedded images
        if ocr and ocr.available and total_chars < max_chars:
            try:
                for rel in doc.part.rels.values():
                    if "image" in rel.reltype:
                        stats.images_found += 1
                        try:
                            image_data = rel.target_part.blob
                            ocr_text = _run_ocr_on_bytes(ocr, image_data)
                            if ocr_text:
                                parts.append(f"[OCR] {ocr_text}")
                                total_chars += len(ocr_text)
                                stats.ocr_pages += 1
                                stats.ocr_chars += len(ocr_text)
                                if total_chars > max_chars:
                                    break
                        except Exception:
                            stats.warnings.append("OCR failed on Word image")
            except Exception:
                pass

        return "\n".join(parts)
    except Exception:
        logger.warning("Failed to extract text from Word document", exc_info=True)
        return ""


def _run_ocr_on_bytes(ocr: OCRService, image_data: bytes) -> str:
    """Run OCR on raw image bytes (synchronous)."""
    import pytesseract
    from PIL import Image

    img = Image.open(io.BytesIO(image_data))
    return str(pytesseract.image_to_string(img)).strip()


def _extract_excel(data: bytes) -> str:
    """Extract text from Excel workbooks using openpyxl."""
    try:
        from openpyxl import load_workbook

        wb = load_workbook(io.BytesIO(data), data_only=True, read_only=True)
        parts: list[str] = []

        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            parts.append(f"=== Sheet: {sheet_name} ===")
            for row in sheet.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                if any(cells):
                    parts.append("\t".join(cells))

        wb.close()
        return "\n".join(parts)
    except Exception:
        logger.warning("Failed to extract text from Excel workbook", exc_info=True)
        return ""


def _extract_powerpoint(data: bytes) -> str:
    """Extract text from PowerPoint presentations using python-pptx."""
    try:
        from pptx import Presentation

        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []

        for i, slide in enumerate(prs.slides, 1):
            slide_texts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_texts.append(para.text)
            if slide_texts:
                parts.append(f"--- Slide {i} ---\n" + "\n".join(slide_texts))

        return "\n\n".join(parts)
    except Exception:
        logger.warning("Failed to extract text from PowerPoint", exc_info=True)
        return ""
